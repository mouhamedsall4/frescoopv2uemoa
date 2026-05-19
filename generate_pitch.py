from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
GREEN = RGBColor(0x1a, 0x7a, 0x5c)
DARK = RGBColor(0x0f, 0x1f, 0x1a)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xFD, 0xF8)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ORANGE = RGBColor(0xE8, 0x6A, 0x17)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_para(tf, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_before = space_before
    return p


# ===== SLIDE 1: TITRE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.6),
         "Hackathon Filières Agricoles GIM-UEMOA 2026 — Point 4 : Accès au financement",
         size=14, color=GREEN, align=PP_ALIGN.LEFT)
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(2),
         "FresCoop", size=60, bold=True, color=WHITE)
tf = add_text(slide, Inches(1), Inches(3.5), Inches(10), Inches(1.5),
              "De l'invisible au finançable. En 90 jours.", size=28, color=WHITE)
add_para(tf, "", size=12)
add_para(tf, "Chaque vente, chaque livraison, chaque paiement d'un agriculteur\ndevient une preuve bancaire vérifiable.", size=16, color=RGBColor(0xA0, 0xD0, 0xC0))


# ===== SLIDE 2: PROBLEME =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.4),
         "PROBLÈME", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(1.2),
         "Les agriculteurs sont invisibles\npour le système financier.", size=36, bold=True, color=DARK)

# 3 stats
stats = [
    ("80%", "des agriculteurs\nsans accès au crédit"),
    ("3%", "du crédit bancaire\nva à l'agriculture"),
    ("0", "historique financier\npour les banques"),
]
for i, (num, desc) in enumerate(stats):
    left = Inches(1 + i * 4)
    add_text(slide, left, Inches(3.2), Inches(3.5), Inches(0.9),
             num, size=52, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    add_text(slide, left, Inches(4.3), Inches(3.5), Inches(1),
             desc, size=16, color=GRAY, align=PP_ALIGN.LEFT)

add_text(slide, Inches(1), Inches(6), Inches(11), Inches(0.8),
         "Pas parce qu'ils ne sont pas fiables — parce qu'ils n'ont aucune preuve exploitable.",
         size=16, color=GRAY, align=PP_ALIGN.LEFT)


# ===== SLIDE 3: SOLUTION =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.4),
         "SOLUTION", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(1),
         "Chaque transaction = une preuve.\nChaque preuve = un score.", size=34, bold=True, color=DARK)

steps = [
    ("1", "Captation automatique", "Ventes, livraisons, paiements → données de scoring"),
    ("2", "Scoring intelligent", "Score 0-100 en temps réel, basé sur l'activité réelle"),
    ("3", "Dossier bancaire portable", "PDF vérifiable par QR code, présentable à toute banque"),
    ("4", "Accès au crédit", "En 90 jours : d'invisible à finançable"),
]
for i, (num, title, desc) in enumerate(steps):
    top = Inches(2.8 + i * 1.1)
    add_text(slide, Inches(1), top, Inches(0.5), Inches(0.5),
             num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Circle behind number
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), top, Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    # Reorder so text on top
    add_text(slide, Inches(1.05), top + Pt(3), Inches(0.4), Inches(0.4),
             num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.7), top, Inches(3.5), Inches(0.4),
             title, size=18, bold=True, color=DARK)
    add_text(slide, Inches(1.7), top + Inches(0.35), Inches(8), Inches(0.4),
             desc, size=14, color=GRAY)


# ===== SLIDE 4: COMMENT CA MARCHE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.4),
         "COMMENT ÇA MARCHE", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         "De l'inscription au premier crédit", size=32, bold=True, color=DARK)

parcours = [
    "Inscription + vérification CNI",
    "Publication des produits",
    "Réception de commandes",
    "Paiement confirmé (PayDunya)",
    "Score qui monte automatiquement",
    "Dossier de crédit exportable",
    "Présentation à une banque/SFD",
    "Obtention du microcrédit",
]
for i, step in enumerate(parcours):
    col = i // 4
    row = i % 4
    left = Inches(1 + col * 6)
    top = Inches(2.3 + row * 1.2)
    add_text(slide, left, top, Inches(0.5), Inches(0.5),
             str(i + 1), size=16, bold=True, color=GREEN)
    add_text(slide, left + Inches(0.6), top, Inches(5), Inches(0.5),
             step, size=17, color=DARK)


# ===== SLIDE 5: GARANTIES =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.4),
         "GARANTIES", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         "4 couches de sécurité. Zéro risque pour la banque.", size=30, bold=True, color=DARK)

garanties = [
    ("Prélèvement à la source", "25% retenu automatiquement à chaque transaction"),
    ("Taux progressif", "0% sous 75K, 15-30% au-dessus — protège le minimum vital"),
    ("Caution solidaire digitale", "Groupes de 5 agriculteurs, 98% de remboursement"),
    ("Fonds de garantie mutualisé", "2% par transaction → couvre 20% des impayés"),
]
for i, (title, desc) in enumerate(garanties):
    left = Inches(1 + (i % 2) * 6)
    top = Inches(2.5 + (i // 2) * 2.2)
    add_text(slide, left, top, Inches(0.4), Inches(0.4),
             str(i + 1), size=14, bold=True, color=GREEN)
    add_text(slide, left + Inches(0.5), top, Inches(5), Inches(0.4),
             title, size=18, bold=True, color=DARK)
    add_text(slide, left + Inches(0.5), top + Inches(0.45), Inches(5.2), Inches(0.7),
             desc, size=14, color=GRAY)

# KPI bar
add_text(slide, Inches(1), Inches(6.3), Inches(3.5), Inches(0.8),
         "Taux de défaut : < 2%", size=16, bold=True, color=GREEN)
add_text(slide, Inches(5), Inches(6.3), Inches(3.5), Inches(0.8),
         "Remboursement : 6-8 mois", size=16, bold=True, color=GREEN)
add_text(slide, Inches(9), Inches(6.3), Inches(3.5), Inches(0.8),
         "Couverture banque : 100%", size=16, bold=True, color=GREEN)


# ===== SLIDE 6: MODELE ECONOMIQUE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.4),
         "MODÈLE ÉCONOMIQUE", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         "Gratuit pour l'agriculteur. Rentable pour tous.", size=30, bold=True, color=DARK)

# Flow: Agriculteur → Acheteur → Banque
flow_items = [("Agriculteur", "GRATUIT"), ("Acheteur", "Paie le produit"), ("SFD / Banque", "Paie l'accès au score")]
for i, (name, role) in enumerate(flow_items):
    left = Inches(1.5 + i * 4)
    add_text(slide, left, Inches(2.3), Inches(3), Inches(0.4),
             name, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(2.7), Inches(3), Inches(0.4),
             role, size=14, color=GREEN, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text(slide, left + Inches(3), Inches(2.3), Inches(0.8), Inches(0.4),
                 "→", size=24, color=GRAY, align=PP_ALIGN.CENTER)

# 3 revenue sources
revenues = [
    ("60%", "Commission marketplace", "2-5% par transaction"),
    ("30%", "Scoring-as-a-Service", "150-500K FCFA/mois par SFD"),
    ("10%", "Services à valeur ajoutée", "Commission 5-15%"),
]
for i, (pct, name, pricing) in enumerate(revenues):
    left = Inches(1 + i * 4)
    top = Inches(3.8)
    add_text(slide, left, top, Inches(3.5), Inches(0.5),
             pct + " du CA", size=14, bold=True, color=GREEN)
    add_text(slide, left, top + Inches(0.45), Inches(3.5), Inches(0.4),
             name, size=17, bold=True, color=DARK)
    add_text(slide, left, top + Inches(0.9), Inches(3.5), Inches(0.4),
             pricing, size=13, color=GRAY)

# Bottom KPIs
kpis = [
    ("Point mort", "100 agri + 3 SFD"),
    ("Coût marginal", "~0 FCFA/agri"),
    ("TAM", "60M × 8 pays"),
    ("Projection an 1", "5 000 utilisateurs"),
]
for i, (label, value) in enumerate(kpis):
    left = Inches(1 + i * 3.1)
    add_text(slide, left, Inches(5.8), Inches(2.8), Inches(0.3),
             label, size=11, color=GRAY)
    add_text(slide, left, Inches(6.1), Inches(2.8), Inches(0.5),
             value, size=16, bold=True, color=DARK)


# ===== SLIDE 7: IMPACT =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(4), Inches(0.4),
         "IMPACT", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         "Des résultats concrets et mesurables.", size=30, bold=True, color=DARK)

impacts = [
    ("65%", "deviennent bancables\nen 90 jours"),
    ("-80%", "coût d'évaluation\npour les SFD"),
    ("90j", "pour obtenir son\npremier crédit"),
    ("×5", "plus d'agriculteurs\nfinancés par SFD"),
]
for i, (num, desc) in enumerate(impacts):
    left = Inches(1 + i * 3.1)
    add_text(slide, left, Inches(2.5), Inches(2.8), Inches(0.8),
             num, size=44, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, left, Inches(3.7), Inches(2.8), Inches(0.8),
             desc, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Bonus points
bonuses = [
    "40% femmes agricultrices ciblées — le scoring élimine les biais de genre",
    "Crédit sans garantie foncière — la preuve d'activité suffit",
    "Souveraineté des données — l'agriculteur contrôle son dossier",
]
for i, bonus in enumerate(bonuses):
    add_text(slide, Inches(1.5), Inches(5 + i * 0.55), Inches(10), Inches(0.5),
             "• " + bonus, size=14, color=DARK)


# ===== SLIDE 8: SCALABILITE UEMOA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.4),
         "SCALABILITÉ", size=12, bold=True, color=GREEN)
add_text(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
         "Un scoring pour 8 pays. 60 millions d'agriculteurs.", size=30, bold=True, color=DARK)

countries = ["Sénégal (pilote)", "Côte d'Ivoire", "Mali", "Burkina Faso", "Bénin", "Niger", "Togo", "Guinée-Bissau"]
for i, country in enumerate(countries):
    col = i // 4
    row = i % 4
    left = Inches(1.5 + col * 5.5)
    top = Inches(2.8 + row * 1)
    marker = "🚀 " if i == 0 else "○ "
    add_text(slide, left, top, Inches(5), Inches(0.5),
             marker + country, size=18, color=DARK if i == 0 else GRAY, bold=(i == 0))

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Infrastructure 100% cloud • Partenariats GIM-UEMOA • Déploiement progressif",
         size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ===== SLIDE 9: DEMO + CTA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1),
         "Ce n'est pas un prototype.", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
         "C'est un produit fonctionnel.", size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
         "Testez le scoring en conditions réelles maintenant.", size=20, color=RGBColor(0xA0, 0xD0, 0xC0), align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6), Inches(11), Inches(0.5),
         "FresCoop — De l'invisible au finançable.", size=16, color=GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = os.path.join(os.path.dirname(__file__), "FRESCOOP_PITCH_v3.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
