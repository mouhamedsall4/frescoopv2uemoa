"""
FresCoop AgriScore - Pitch Deck GIM-UEMOA Hackathon 2026
Theme: Acces limite au financement agricole
6 slides: Accroche > Probleme > Solution > Modele Economique > Impact > CTA
Version finale avec modele economique complet
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Theme colors
GREEN_DARK = RGBColor(0x0A, 0x4D, 0x3A)
GREEN_MID = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GOLD = RGBColor(0xD9, 0x77, 0x06)
GOLD_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
CYAN = RGBColor(0x08, 0x91, 0xB2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_shape_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16, color=BLACK, align=PP_ALIGN.LEFT, bold=False, spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.alignment = align
        p.space_after = Pt(font_size * (spacing - 1))
    return txBox


def add_icon_bullet(slide, left, top, width, height, icon_text, title, desc, accent_color=GREEN_MID):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.1), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]
    p_c.text = icon_text
    p_c.font.size = Pt(13)
    p_c.font.color.rgb = WHITE
    p_c.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.6), top, width - Inches(0.6), Inches(0.3),
                 title, font_size=14, bold=True, color=BLACK)
    add_text_box(slide, left + Inches(0.6), top + Inches(0.32), width - Inches(0.6), height - Inches(0.32),
                 desc, font_size=11, color=GRAY)


# ============================================================================
# SLIDE 1: ACCROCHE / TITRE
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide1, GREEN_DARK)

# Logo
logo = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(0.5), Inches(0.9), Inches(0.9))
logo.fill.solid()
logo.fill.fore_color.rgb = RGBColor(0x15, 0x80, 0x3D)
logo.line.color.rgb = GREEN_MID
logo.line.width = Pt(3)
tf_logo = logo.text_frame
tf_logo.vertical_anchor = MSO_ANCHOR.MIDDLE
p_logo = tf_logo.paragraphs[0]
p_logo.text = "F"
p_logo.font.size = Pt(32)
p_logo.font.bold = True
p_logo.font.color.rgb = GREEN_MID
p_logo.alignment = PP_ALIGN.CENTER

# Title
add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(8), Inches(1.0),
             "FresCoop AgriScore", font_size=48, bold=True, color=WHITE)

# Tagline
add_text_box(slide1, Inches(0.8), Inches(2.8), Inches(8), Inches(0.6),
             "Rendre chaque agriculteur finançable grâce aux données", font_size=22, color=RGBColor(0xBB, 0xF7, 0xD0))

# Accroche
add_multiline_box(slide1, Inches(0.8), Inches(3.7), Inches(7.5), Inches(2.0), [
    "L'agriculture = 35% du PIB UEMOA, 60% des emplois.",
    "Pourtant : moins de 3% du financement formel.",
    "",
    "Pourquoi ? Les banques ne prêtent pas à ceux",
    "qu'elles ne peuvent pas scorer.",
    "",
    "FresCoop résout ce paradoxe."
], font_size=14, color=RGBColor(0xD1, 0xFA, 0xE5))

# Badge hackathon
badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.4), Inches(4.2), Inches(0.6))
badge.fill.solid()
badge.fill.fore_color.rgb = GOLD
badge.line.fill.background()
tf_badge = badge.text_frame
tf_badge.vertical_anchor = MSO_ANCHOR.MIDDLE
p_badge = tf_badge.paragraphs[0]
p_badge.text = "  HACKATHON GIM-UEMOA 2026 — Filières Agricoles"
p_badge.font.size = Pt(12)
p_badge.font.bold = True
p_badge.font.color.rgb = BLACK
p_badge.alignment = PP_ALIGN.LEFT

# Right - big stat box
stat_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.8), Inches(3.5), Inches(2.2))
stat_box.fill.solid()
stat_box.fill.fore_color.rgb = RGBColor(0x15, 0x80, 0x3D)
stat_box.line.fill.background()
tf_stat = stat_box.text_frame
tf_stat.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_stat.word_wrap = True
p1 = tf_stat.paragraphs[0]
p1.text = "$117 Mds"
p1.font.size = Pt(38)
p1.font.bold = True
p1.font.color.rgb = GOLD
p1.alignment = PP_ALIGN.CENTER
p2 = tf_stat.add_paragraph()
p2.text = "Gap de financement"
p2.font.size = Pt(12)
p2.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
p2.alignment = PP_ALIGN.CENTER
p3 = tf_stat.add_paragraph()
p3.text = "agricole en Afrique (IFC)"
p3.font.size = Pt(12)
p3.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
p3.alignment = PP_ALIGN.CENTER

# Right - second stat
stat_box2 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(4.3), Inches(3.5), Inches(1.8))
stat_box2.fill.solid()
stat_box2.fill.fore_color.rgb = RGBColor(0x15, 0x80, 0x3D)
stat_box2.line.fill.background()
tf_s2 = stat_box2.text_frame
tf_s2.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_s2.word_wrap = True
p_s2 = tf_s2.paragraphs[0]
p_s2.text = "97%"
p_s2.font.size = Pt(38)
p_s2.font.bold = True
p_s2.font.color.rgb = RED
p_s2.alignment = PP_ALIGN.CENTER
p_s2b = tf_s2.add_paragraph()
p_s2b.text = "des agriculteurs UEMOA"
p_s2b.font.size = Pt(12)
p_s2b.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
p_s2b.alignment = PP_ALIGN.CENTER
p_s2c = tf_s2.add_paragraph()
p_s2c.text = "sans accès au crédit formel"
p_s2c.font.size = Pt(12)
p_s2c.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
p_s2c.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 2: PROBLEME
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide2, WHITE)

tag2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(1.8), Inches(0.4))
tag2.fill.solid()
tag2.fill.fore_color.rgb = RGBColor(0xFE, 0xE2, 0xE2)
tag2.line.fill.background()
tf_tag = tag2.text_frame
tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
p_tag = tf_tag.paragraphs[0]
p_tag.text = "LE PROBLÈME"
p_tag.font.size = Pt(10)
p_tag.font.bold = True
p_tag.font.color.rgb = RED
p_tag.alignment = PP_ALIGN.CENTER

add_text_box(slide2, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "L'accès au financement agricole est bloqué par un cercle vicieux", font_size=26, bold=True, color=BLACK)

add_text_box(slide2, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
             "L'absence de données fiables empêche les institutions de financer ceux qui produisent la richesse du continent.", font_size=13, color=GRAY)

# 3 problem columns: Agriculteur / Banque / Economie
problems = [
    ("L'AGRICULTEUR", RED, [
        ("Invisible", "Aucun historique bancaire ni scoring"),
        ("Isolé", "Pas de preuve économique de son activité"),
        ("Bloqué", "Ne peut pas investir pour augmenter sa production"),
        ("Méfiant", "Les agents terrain coûtent cher et sont rares"),
    ]),
    ("LA BANQUE / SFD", GOLD, [
        ("Aveugle", "Pas de données pour évaluer le risque"),
        ("Coûts élevés", "Agents terrain = 15-25% du coût du crédit"),
        ("Défauts", "Taux de défaut de 15-30% sans scoring"),
        ("Opportunité manquée", "$117 Mds de marché non servi"),
    ]),
    ("L'ÉCONOMIE UEMOA", BLUE, [
        ("Sous-production", "Agriculteurs sans capital = rendements faibles"),
        ("Exode rural", "Les jeunes quittent faute de perspectives"),
        ("Dépendance", "Importations alimentaires en hausse"),
        ("Gaspillage", "30-40% de pertes post-récolte faute de moyens"),
    ]),
]

for col, (header, accent, items) in enumerate(problems):
    x = Inches(0.8) + Inches(4.1) * col
    y_start = Inches(2.2)
    col_w = Inches(3.8)

    # Header
    hdr = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, col_w, Inches(0.45))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = accent
    hdr.line.fill.background()
    tf_h = hdr.text_frame
    tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_h = tf_h.paragraphs[0]
    p_h.text = header
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = WHITE
    p_h.alignment = PP_ALIGN.CENTER

    for i, (title, desc) in enumerate(items):
        y = y_start + Inches(0.6) + Inches(1.0) * i
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = GRAY_LIGHT
        card.line.fill.background()

        add_text_box(slide2, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.4), Inches(0.3),
                     title, font_size=12, bold=True, color=accent)
        add_text_box(slide2, x + Inches(0.2), y + Inches(0.42), col_w - Inches(0.4), Inches(0.4),
                     desc, font_size=11, color=GRAY)

# Bottom quote
add_text_box(slide2, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
             "« Ce n'est pas un manque de confiance. C'est un manque de données. FresCoop comble ce vide. »",
             font_size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 3: SOLUTION
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide3, WHITE)

tag3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.4))
tag3.fill.solid()
tag3.fill.fore_color.rgb = GREEN_LIGHT
tag3.line.fill.background()
tf_tag3 = tag3.text_frame
tf_tag3.vertical_anchor = MSO_ANCHOR.MIDDLE
p_tag3 = tf_tag3.paragraphs[0]
p_tag3.text = "NOTRE SOLUTION"
p_tag3.font.size = Pt(10)
p_tag3.font.bold = True
p_tag3.font.color.rgb = GREEN_DARK
p_tag3.alignment = PP_ALIGN.CENTER

add_text_box(slide3, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
             "FresCoop : la plateforme qui transforme l'activité agricole en score bancaire", font_size=24, bold=True, color=BLACK)

add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(10), Inches(0.4),
             "Un écosystème complet : marketplace + scoring + crédit + remboursement automatique via GIM-Pay",
             font_size=13, color=GRAY)

# Left: 5 pillars
pillars = [
    ("1", "Marketplace agricole + traçabilité", "Commandes horodatées, livraisons suivies, preuves économiques automatiques"),
    ("2", "Score de bancabilité sur 100 pts", "12 critères : ventes, livraisons, régularité, identité, avis clients (étoiles)"),
    ("3", "Dossier bancaire digital", "Export PDF/API vers banques, SFD et fintechs partenaires"),
    ("4", "Crédit par tranches sécurisées", "40/30/30% — chaque tranche conditionnée à la preuve d'utilisation"),
    ("5", "Remboursement intégré via GIM-Pay", "Prélèvement auto sur chaque vente. Taux de défaut < 5%"),
]

for i, (num, title, desc) in enumerate(pillars):
    y = Inches(2.2) + Inches(1.0) * i
    add_icon_bullet(slide3, Inches(0.8), y, Inches(6.8), Inches(0.85), num, title, desc, GREEN_MID)

# Right: Diagram - scoring visual
diag_title = add_text_box(slide3, Inches(8.2), Inches(2.0), Inches(4.5), Inches(0.4),
             "COMMENT ÇA MARCHE", font_size=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

flow_steps = [
    ("Inscription", "Profil + CNI vérifiée", GREEN_LIGHT, GREEN_DARK),
    ("Ventes", "Commandes sur FresCoop", GOLD_LIGHT, GOLD),
    ("Notation", "Clients notent (étoiles)", GRAY_LIGHT, BLACK),
    ("Score", "Bancabilité calculée /100", GREEN_LIGHT, GREEN_DARK),
    ("Crédit", "Demande en 1 clic", RGBColor(0xCF, 0xFA, 0xFE), CYAN),
    ("Rembourse", "Auto via ventes GIM-Pay", GOLD_LIGHT, GOLD),
]

for i, (step, detail, bg_col, txt_col) in enumerate(flow_steps):
    y = Inches(2.5) + Inches(0.75) * i
    step_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), y, Inches(4.2), Inches(0.6))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = bg_col
    step_box.line.fill.background()

    tf_step = step_box.text_frame
    tf_step.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_step = tf_step.paragraphs[0]
    p_step.text = f"  {step}  →  {detail}"
    p_step.font.size = Pt(11)
    p_step.font.bold = True
    p_step.font.color.rgb = txt_col
    p_step.font.name = 'Segoe UI'

    if i < len(flow_steps) - 1:
        arrow = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.3), y + Inches(0.55), Inches(0.18), Inches(0.18))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREEN_MID
        arrow.line.fill.background()


# ============================================================================
# SLIDE 4: MODELE ECONOMIQUE (COMPLET)
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide4, WHITE)

tag4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(3.0), Inches(0.4))
tag4.fill.solid()
tag4.fill.fore_color.rgb = GOLD_LIGHT
tag4.line.fill.background()
tf_tag4 = tag4.text_frame
tf_tag4.vertical_anchor = MSO_ANCHOR.MIDDLE
p_tag4 = tf_tag4.paragraphs[0]
p_tag4.text = "MODÈLE ÉCONOMIQUE"
p_tag4.font.size = Pt(10)
p_tag4.font.bold = True
p_tag4.font.color.rgb = GOLD
p_tag4.alignment = PP_ALIGN.CENTER

add_text_box(slide4, Inches(0.8), Inches(0.95), Inches(11), Inches(0.5),
             "Un modèle gagnant-gagnant-gagnant : tout le monde y gagne", font_size=24, bold=True, color=BLACK)

add_text_box(slide4, Inches(0.8), Inches(1.45), Inches(11), Inches(0.3),
             "FresCoop monétise la création de valeur, pas la pauvreté. Chaque acteur paye pour ce qu'il reçoit.",
             font_size=12, color=GRAY)

# 4 columns: Agriculteur / Banque-SFD / Client / Agents
cols = [
    ("AGRICULTEUR", GREEN_MID, [
        "CE QU'IL REÇOIT :",
        "• Visibilité marketplace",
        "• Accès au financement",
        "• Score de bancabilité",
        "• Écoulement garanti",
        "",
        "CE QU'IL PAYE :",
        "• 3-5% sur ses ventes",
        "  (0% les 50K premiers FCFA)",
        "• Invisible dans son gain",
        "  (il vend 15-30% plus cher)",
    ]),
    ("BANQUE / SFD", BLUE, [
        "CE QU'ELLE REÇOIT :",
        "• Clients pré-scorés",
        "• Scoring fiable (vs agents)",
        "• Coût 10x < agents terrain",
        "• Taux défaut < 5%",
        "",
        "CE QU'ELLE PAYE :",
        "• Abo 100K-200K FCFA/mois",
        "  (scoring illimité)",
        "• 1.5-2.5% sur prêts débloqués",
        "• 15% du revenu d'intérêt",
    ]),
    ("CLIENT / ACHETEUR", CYAN, [
        "CE QU'IL REÇOIT :",
        "• Produits frais traçables",
        "• Prix justes transparents",
        "• Cashback 2% fidélité",
        "• Notation → confiance",
        "",
        "CE QU'IL PAYE :",
        "• Prix normal du produit",
        "• (Pas de frais supplémentaire)",
        "",
        "→ Le cashback le fidélise",
        "→ Ses avis scorent les agris",
    ]),
    ("AGENTS TERRAIN", GOLD, [
        "CE QU'ILS REÇOIVENT :",
        "• Revenus complémentaires",
        "• App mobile dédiée",
        "• Formation certifiante",
        "",
        "CE QU'ILS GAGNENT :",
        "• 1500 FCFA/agri inscrit",
        "• 1.5% sur prêts facilités",
        "• Bonus 15K si rembours >95%",
        "",
        "PARTENAIRES : ANCAR, SAED",
        "(700+ agents déjà déployés)",
    ]),
]

col_w = Inches(2.9)
col_gap = Inches(0.2)

for i, (header, accent, lines) in enumerate(cols):
    x = Inches(0.6) + (col_w + col_gap) * i
    y_start = Inches(1.9)

    # Header
    hdr = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, col_w, Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = accent
    hdr.line.fill.background()
    tf_h = hdr.text_frame
    tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_h = tf_h.paragraphs[0]
    p_h.text = header
    p_h.font.size = Pt(10)
    p_h.font.bold = True
    p_h.font.color.rgb = WHITE
    p_h.alignment = PP_ALIGN.CENTER

    # Content card
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start + Inches(0.45), col_w, Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = GRAY_LIGHT
    card.line.fill.background()

    txBox = slide4.shapes.add_textbox(x + Inches(0.15), y_start + Inches(0.55), col_w - Inches(0.3), Inches(3.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        is_header_line = line.startswith("CE QU") or line.startswith("PARTENAIRES") or line.startswith("→")
        p.font.size = Pt(9.5 if not is_header_line else 9)
        p.font.bold = is_header_line
        p.font.color.rgb = accent if is_header_line else BLACK
        p.font.name = 'Segoe UI'
        p.space_after = Pt(2)

# Bottom: Revenue projection
rev_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.9))
rev_box.fill.solid()
rev_box.fill.fore_color.rgb = GREEN_DARK
rev_box.line.fill.background()

rev_stats = [
    ("3-5%", "Commission\nagriculteur"),
    ("1.5-2.5%", "Fee prêt\ndébloqué"),
    ("100-200K/mois", "Abo scoring\nbanque/SFD"),
    ("2%", "Cashback client\n(fidélisation)"),
    ("33M", "Agriculteurs\nadressables"),
]

for i, (num, label) in enumerate(rev_stats):
    x = Inches(0.9) + Inches(2.5) * i

    add_text_box(slide4, x, Inches(6.38), Inches(2.2), Inches(0.35),
                 num, font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txBox_m = slide4.shapes.add_textbox(x, Inches(6.72), Inches(2.2), Inches(0.5))
    tf_m = txBox_m.text_frame
    tf_m.word_wrap = True
    for j, line in enumerate(label.split('\n')):
        if j == 0:
            p = tf_m.paragraphs[0]
        else:
            p = tf_m.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
        p.font.name = 'Segoe UI'
        p.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 5: IMPACT
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide5, WHITE)

tag5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.4))
tag5.fill.solid()
tag5.fill.fore_color.rgb = GREEN_LIGHT
tag5.line.fill.background()
tf_tag5 = tag5.text_frame
tf_tag5.vertical_anchor = MSO_ANCHOR.MIDDLE
p_tag5 = tf_tag5.paragraphs[0]
p_tag5.text = "IMPACT"
p_tag5.font.size = Pt(10)
p_tag5.font.bold = True
p_tag5.font.color.rgb = GREEN_DARK
p_tag5.alignment = PP_ALIGN.CENTER

add_text_box(slide5, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "Impact social, économique et systémique mesurable", font_size=24, bold=True, color=BLACK)

add_text_box(slide5, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
             "Chaque agriculteur financé = une famille qui sort de la précarité + une chaîne de valeur structurée + un client bancaire créé.",
             font_size=12, color=GRAY)

# 2x3 grid
impacts = [
    ("+40%", "Inclusion financière", "d'agriculteurs accédant au crédit\nformel pour la première fois", GREEN_MID),
    ("< 5%", "Taux de défaut", "vs 15-30% sans scoring\n(remboursement auto sur ventes)", BLUE),
    ("+60%", "Revenus agriculteurs", "grâce à l'accès au capital\net à la marketplace directe", GOLD),
    ("÷10", "Coût d'évaluation", "pour les banques vs agents terrain\n(scoring digital < inspection)", GREEN_MID),
    ("-35%", "Pertes post-récolte", "les agriculteurs financés stockent\net vendent au bon moment", BLUE),
    ("x3", "Emplois indirects", "par chaîne de valeur financée\n(transport, stockage, commerce)", GOLD),
]

for i, (metric, title, desc, color) in enumerate(impacts):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + Inches(4.1) * col
    y = Inches(2.1) + Inches(2.2) * row

    card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.9))
    card.fill.solid()
    card.fill.fore_color.rgb = GRAY_LIGHT
    card.line.fill.background()

    add_text_box(slide5, x + Inches(0.25), y + Inches(0.15), Inches(1.5), Inches(0.5),
                 metric, font_size=28, bold=True, color=color)

    add_text_box(slide5, x + Inches(1.8), y + Inches(0.2), Inches(1.8), Inches(0.35),
                 title, font_size=12, bold=True, color=BLACK)

    txBox_d = slide5.shapes.add_textbox(x + Inches(0.25), y + Inches(0.85), Inches(3.3), Inches(0.9))
    tf_d = txBox_d.text_frame
    tf_d.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            p = tf_d.paragraphs[0]
        else:
            p = tf_d.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.font.name = 'Segoe UI'

# ODD + Partenariats
add_text_box(slide5, Inches(0.8), Inches(6.5), Inches(6), Inches(0.7),
             "ODD : 1 (Pauvreté) • 2 (Faim) • 8 (Emploi) • 9 (Innovation) • 10 (Inégalités)",
             font_size=10, color=GRAY)

add_text_box(slide5, Inches(7.0), Inches(6.5), Inches(5.5), Inches(0.7),
             "Partenaires stratégiques : GIM-UEMOA • ANCAR • La Banque Agricole • SAED • Wave",
             font_size=10, bold=True, color=GREEN_DARK)


# ============================================================================
# SLIDE 6: CALL TO ACTION
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide6, GREEN_DARK)

add_text_box(slide6, Inches(0.8), Inches(1.0), Inches(8), Inches(0.8),
             "FresCoop AgriScore", font_size=42, bold=True, color=WHITE)

add_text_box(slide6, Inches(0.8), Inches(1.8), Inches(8), Inches(0.6),
             "De l'invisible au finançable.", font_size=26, color=GREEN_MID)

# Differentiators
diff_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.7), Inches(7.0), Inches(3.0))
diff_box.fill.solid()
diff_box.fill.fore_color.rgb = RGBColor(0x15, 0x80, 0x3D)
diff_box.line.fill.background()

diffs = [
    "Seul scoring agricole natif UEMOA (pas un copier-coller kenyan)",
    "Intégration GIM-Pay native — paiement dans l'écosystème du jury",
    "Remboursement zéro-défaut intégré aux ventes",
    "Réseau d'agents existant (ANCAR 700+ / SAED) — pas à construire",
    "Déjà fonctionnel : Web + Mobile + API + Mode offline",
    "3 filières GIM : Arachide • Mangue • Anacarde",
    "Cashback client 2% — cercle vertueux de données",
]

txBox_diff = slide6.shapes.add_textbox(Inches(1.1), Inches(2.85), Inches(6.5), Inches(2.8))
tf_diff = txBox_diff.text_frame
tf_diff.word_wrap = True
for i, d in enumerate(diffs):
    if i == 0:
        p = tf_diff.paragraphs[0]
    else:
        p = tf_diff.add_paragraph()
    p.text = f"  ✓  {d}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    p.font.name = 'Segoe UI'
    p.space_after = Pt(5)

# Right: Notre demande
ask_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(2.7), Inches(4.3), Inches(3.0))
ask_box.fill.solid()
ask_box.fill.fore_color.rgb = GOLD
ask_box.line.fill.background()

tf_ask = ask_box.text_frame
tf_ask.vertical_anchor = MSO_ANCHOR.MIDDLE
tf_ask.word_wrap = True

asks = [
    ("NOTRE DEMANDE", 11, True, BLACK),
    ("", 6, False, BLACK),
    ("1. Accès API GIM-Pay", 14, True, BLACK),
    ("   Paiements et scoring intégrés", 11, False, RGBColor(0x45, 0x30, 0x03)),
    ("", 6, False, BLACK),
    ("2. Pilote avec 3 SFD partenaires", 14, True, BLACK),
    ("   sur Arachide + Mangue (6 mois)", 11, False, RGBColor(0x45, 0x30, 0x03)),
    ("", 6, False, BLACK),
    ("3. Convention ANCAR", 14, True, BLACK),
    ("   Accès au réseau de 700 agents", 11, False, RGBColor(0x45, 0x30, 0x03)),
]

for i, (text, size, bold, color) in enumerate(asks):
    if i == 0:
        p = tf_ask.paragraphs[0]
    else:
        p = tf_ask.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT

# Bottom tagline
add_text_box(slide6, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             "« Quand l'Afrique finance ses propres agriculteurs, tout le continent s'élève. »",
             font_size=14, bold=True, color=RGBColor(0x86, 0xEF, 0xAC), align=PP_ALIGN.CENTER)

add_text_box(slide6, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
             "Hackathon GIM-UEMOA 2026  •  Filières Agricoles  •  Équipe FresCoop  •  frescoop.vercel.app",
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# Save
output_path = r'C:\Users\HP\frescoopuemoa-v2\FRESCOOP_PITCH_HACKATHON_GIMUEMOA.pptx'
prs.save(output_path)
print(f"Pitch deck genere: {output_path}")
print("6 slides finales:")
print("  1. Accroche ($117Mds gap + 97% exclus)")
print("  2. Probleme (Agriculteur / Banque / Economie)")
print("  3. Solution (5 piliers + parcours)")
print("  4. Modele Economique (4 acteurs: agri/banque/client/agents)")
print("  5. Impact (6 metriques + ODD + partenaires)")
print("  6. CTA (7 differenciants + 3 demandes)")
