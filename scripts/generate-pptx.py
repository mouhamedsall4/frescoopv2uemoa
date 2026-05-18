"""
FresCoop — Pitch Deck PPTX Generator
Generates a professional 12-slide presentation for the GIM-UEMOA Hackathon
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# Colors
GREEN_DARK = RGBColor(0x06, 0x4E, 0x3B)
GREEN_MID = RGBColor(0x0F, 0x76, 0x6E)
GREEN_LIGHT = RGBColor(0x10, 0xB9, 0x81)
GREEN_PALE = RGBColor(0xD1, 0xFA, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xD9, 0x99, 0x12)
CORAL = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x25, 0x83, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, color1, color2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=14, color=DARK, bold_first=False, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox

def add_shape_card(slide, left, top, width, height, fill_color, text='', font_size=12, font_color=DARK, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Segoe UI'
    return shape

def add_kpi_card(slide, left, top, value, label, color=GREEN_DARK):
    shape = add_shape_card(slide, left, top, Inches(2.6), Inches(1.5), color)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Segoe UI'

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=DARK, bullet_color=GREEN_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 : COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, GREEN_DARK, GREEN_MID)

# Badge
add_shape_card(slide, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.5), RGBColor(0x13, 0x6B, 0x5A), 'HACKATHON FILIERES AGRICOLES GIM-UEMOA 2026', 9, WHITE)

# Title
add_text_box(slide, Inches(2), Inches(2.3), Inches(9), Inches(1.5), 'FresCoop', 60, True, WHITE, PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2.5), Inches(3.6), Inches(8), Inches(0.6), 'De la recolte au credit : la chaine de valeur agricole digitalisee', 22, False, GREEN_PALE, PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(2.5), Inches(4.5), Inches(8), Inches(1), 'Micro-hubs solaires, intelligence marche, scoring agricole (AgriScore)\net paiement partenaire pour rendre chaque agriculteur financable', 14, False, RGBColor(0xA7, 0xF3, 0xD0), PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, Inches(3), Inches(6.3), Inches(7), Inches(0.5), 'Friction 4 : Acces limite aux financements agricoles  |  Equipe FresCoop  |  Mai 2026', 10, False, RGBColor(0x6E, 0xE7, 0xB7), PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 : LE PROBLEME
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'LE PROBLEME', 12, True, GREEN_MID, PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8), '80% des agriculteurs UEMOA sont exclus du credit.\nPas parce qu\'ils sont risques — parce que leur credibilite n\'est pas documentee.', 20, True, DARK, PP_ALIGN.LEFT)

# KPI cards
add_kpi_card(slide, Inches(0.8), Inches(2.2), '80%', 'Agriculteurs exclus\ndu credit formel', GREEN_DARK)
add_kpi_card(slide, Inches(3.6), Inches(2.2), '3%', 'Credit bancaire vers\nl\'agriculture', CORAL)
add_kpi_card(slide, Inches(6.4), Inches(2.2), '60M', 'Exploitants agricoles\ndans l\'espace UEMOA', BLUE)
add_kpi_card(slide, Inches(9.2), Inches(2.2), '$2.5Mds', 'Gap de financement\nagricole Afrique Ouest', GOLD)

# Explication
add_multiline_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(3), [
    'Pourquoi ce probleme persiste :',
    '',
    '   Les banques n\'ont aucune donnee fiable sur les agriculteurs',
    '   Les ventes existent mais sont informelles et non tracees',
    '   Le titre foncier formel exclut 70% des producteurs',
    '   Les agriculteurs productifs remboursent deja (tontines, prets informels)',
    '   Mais cette credibilite n\'est documentee nulle part',
], 13, DARK, True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 : LA SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'NOTRE SOLUTION', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'FresCoop : du terrain au dossier bancaire en 5 etapes', 22, True, DARK)

# 5 piliers en cards
piliers = [
    ('1', 'Collecte\nassistee', 'Agent terrain\nLangue locale\nTemoin + consentement'),
    ('2', 'Score\ncontextualise', '13 criteres\nFoncier coutumier\nTontine, Mobile Money'),
    ('3', 'Dossier\nverifiable', 'PDF + QR code\nPreuves tracees\nCode verification'),
    ('4', 'Calendrier\nsaisonnier', 'Remboursement aligne\nsur cycles recolte\net periodes vente'),
    ('5', 'Boucle\npost-credit', 'Suivi remboursement\nHistorique evolutif\nScore qui progresse'),
]

for i, (num, title, desc) in enumerate(piliers):
    left = Inches(0.5 + i * 2.5)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), Inches(1.8), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_DARK
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, left, Inches(2.6), Inches(2.3), Inches(0.7), title, 13, True, GREEN_DARK, PP_ALIGN.CENTER)

    # Description card
    card = add_shape_card(slide, left, Inches(3.4), Inches(2.3), Inches(1.8), GREEN_PALE, desc, 10, DARK)

# Arrow flow at bottom
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5), 'AGRICULTEUR  ───────>  AGENT TERRAIN  ───────>  FRESCOOP (AgriScore)  ───────>  BANQUE / IMF  ───────>  CREDIT', 11, True, GREEN_MID, PP_ALIGN.CENTER)

# Key message
add_shape_card(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.8), GREEN_DARK, 'L\'agriculteur ne s\'inscrit pas seul sur une app. Il est accompagne par un agent de confiance, dans sa langue.', 12, WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 : SCORING — CE QUI NOUS DIFFERENCIE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'SCORING CONTEXTUALISE', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), '13 criteres adaptes a la realite agricole africaine', 20, True, DARK)

# Two columns comparison
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), 'FresCoop — Scoring AgriScore', 14, True, GREEN_DARK)
criteria_left = [
    'Anciennete du compte agriculteur',
    'Transactions verifiees (ventes reelles)',
    'Paiements Mobile Money (Wave, Orange)',
    'Lots traces avec capteurs IoT',
    'Attestations de cooperatives',
    'Preuves economiques documentees',
    'Revenu mensuel moyen calcule',
    'Groupement / GIE / Cooperative',
    'Foncier (formel OU coutumier)',
    'Experience agricole (annees)',
    'Mobile Money actif',
    'Tontine / epargne informelle',
    'Historique de remboursement',
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(5), [f'  {c}' for c in criteria_left], 11, DARK)

add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4), 'Scoring bancaire classique', 14, True, CORAL)
criteria_right = [
    'Titre foncier formel obligatoire',
    'Bulletin de salaire requis',
    'Historique bancaire (compte courant)',
    'Garanties immobilieres',
    'Pas de prise en compte du coutumier',
    'Pas de donnees de production',
    'Pas de preuves informelles',
    'Exclut 80% des agriculteurs',
]
add_bullet_list(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4), [f'  {c}' for c in criteria_right], 11, GRAY)

# Bottom highlight
add_shape_card(slide, Inches(7), Inches(5.8), Inches(5.5), Inches(1.2), RGBColor(0xFE, 0xE2, 0xE2), 'Resultat : 80% des agriculteurs\nexclus du systeme', 12, CORAL)
add_shape_card(slide, Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.2), GREEN_PALE, 'Resultat : scoring inclusif\nqui valorise les realites terrain', 12, GREEN_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 : DEMO / PARCOURS UTILISATEUR
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'PARCOURS UTILISATEUR', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'De l\'inscription au premier credit : le parcours reel', 20, True, DARK)

steps = [
    ('1', 'INSCRIPTION', 'L\'agriculteur s\'inscrit\n(ou l\'agent l\'inscrit)\nen 2 minutes'),
    ('2', 'COLLECTE', 'L\'agent collecte les\ndonnees terrain en\nlangue locale'),
    ('3', 'SCORING', 'Le score se calcule\nautomatiquement\n(13 criteres)'),
    ('4', 'DOSSIER', 'Un dossier bancaire\nPDF est genere avec\nQR de verification'),
    ('5', 'CREDIT', 'L\'IMF consulte le\ndossier et decide\n(5 min vs 3 semaines)'),
    ('6', 'SUIVI', 'Remboursements suivis\nLe score evolue\nHistorique construit'),
]

for i, (num, title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.8 + row * 2.7)

    # Card
    card = add_shape_card(slide, left, top, Inches(3.7), Inches(2.2), GREEN_PALE if row == 0 else RGBColor(0xEC, 0xFD, 0xF5))

    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1), top + Inches(0.1), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_DARK
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), Inches(2.8), Inches(0.4), title, 13, True, GREEN_DARK)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.3), Inches(1.4), desc, 11, False, DARK, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 : BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'BUSINESS MODEL', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'L\'agriculteur ne paie jamais. Les IMF paient pour des clients pre-qualifies.', 18, True, DARK)

# Revenue sources
sources = [
    ('45%', 'Dossiers qualifies', '2 000 - 5 000 FCFA\npar dossier consulte', GREEN_DARK),
    ('30%', 'Abonnements IMF', '500K - 2M FCFA/mois\nlicence plateforme + API', GREEN_MID),
    ('15%', 'Commission credit', '1-2% du montant\ndu pret debloque', BLUE),
    ('10%', 'Programmes', 'Contrats bailleurs\net gouvernements', GOLD),
]

for i, (pct, title, desc, color) in enumerate(sources):
    left = Inches(0.5 + i * 3.15)
    top = Inches(1.8)

    shape = add_shape_card(slide, left, top, Inches(2.9), Inches(2.5), color)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = pct
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = ''
    p3.font.size = Pt(6)
    p4 = tf.add_paragraph()
    p4.text = desc
    p4.font.size = Pt(10)
    p4.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    p4.alignment = PP_ALIGN.CENTER

# Revenue projection
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(5), Inches(0.4), 'Projection de revenus (FCFA)', 14, True, DARK)

# Chart
chart_data = CategoryChartData()
chart_data.categories = ['Annee 1\n(Pilote)', 'Annee 2\n(Croissance)', 'Annee 3\n(Echelle)']
chart_data.add_series('Revenus', (46, 201, 480))
chart_data.add_series('Couts', (46, 149, 282))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.8), Inches(5.1), Inches(6), Inches(2.2),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(9)
plot = chart.plots[0]
series_rev = plot.series[0]
series_rev.format.fill.solid()
series_rev.format.fill.fore_color.rgb = GREEN_DARK
series_cost = plot.series[1]
series_cost.format.fill.solid()
series_cost.format.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

# Right side KPIs
add_kpi_card(slide, Inches(7.5), Inches(5.0), '480M', 'Revenus Annee 3\n(FCFA)', GREEN_DARK)
add_kpi_card(slide, Inches(10.3), Inches(5.0), '41%', 'Marge nette\nAnnee 3', GREEN_MID)
add_kpi_card(slide, Inches(7.5), Inches(6.0+0.6), 'x10', 'Croissance\nen 3 ans', BLUE)
add_kpi_card(slide, Inches(10.3), Inches(6.0+0.6), '15M', 'Investissement\ninitial (FCFA)', GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 : MARCHE (TAM SAM SOM)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'OPPORTUNITE DE MARCHE', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'Un marche de $2.5 milliards en Afrique de l\'Ouest', 20, True, DARK)

# Concentric circles TAM SAM SOM
# TAM
tam = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(1.8), Inches(5), Inches(5))
tam.fill.solid()
tam.fill.fore_color.rgb = GREEN_PALE
tam.line.color.rgb = GREEN_LIGHT
tf = tam.text_frame
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = 'TAM : $2.5 Mds'
p.font.size = Pt(11)
p.font.color.rgb = GREEN_MID
p.alignment = PP_ALIGN.CENTER

# SAM
sam = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(2.8), Inches(3.2), Inches(3.2))
sam.fill.solid()
sam.fill.fore_color.rgb = RGBColor(0xA7, 0xF3, 0xD0)
sam.line.color.rgb = GREEN_MID
tf = sam.text_frame
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = 'SAM : $450M'
p.font.size = Pt(11)
p.font.color.rgb = GREEN_DARK
p.alignment = PP_ALIGN.CENTER

# SOM
som = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(3.8), Inches(1.8), Inches(1.8))
som.fill.solid()
som.fill.fore_color.rgb = GREEN_DARK
som.line.fill.background()
tf = som.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = 'SOM\n$12M'
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Descriptions on the right
descs = [
    ('TAM — $2.5 Milliards/an', 'Scoring credit agricole pour 60M\nd\'exploitants en Afrique de l\'Ouest (8 pays UEMOA)'),
    ('SAM — $450 Millions/an', 'Senegal + Mali + Cote d\'Ivoire\nFilieres porteuses : riz, arachide, mangue, oignon'),
    ('SOM — $12 Millions/an', 'Pilote Senegal : 3 regions, 5 000 agriculteurs,\n10 IMF partenaires (horizon 3 ans)'),
]
for i, (title, desc) in enumerate(descs):
    top = Inches(2.0 + i * 1.7)
    add_text_box(slide, Inches(7.2), top, Inches(5.5), Inches(0.35), title, 13, True, GREEN_DARK)
    add_text_box(slide, Inches(7.2), top + Inches(0.4), Inches(5.5), Inches(0.9), desc, 11, False, GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 : IMPACT SOCIAL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'IMPACT SOCIAL & ENVIRONNEMENTAL', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'Des resultats mesurables a 3 ans', 20, True, DARK)

# Impact KPIs
impacts = [
    ('25 000', 'Agriculteurs avec\nidentite financiere', GREEN_DARK),
    ('$3.4M', 'Credits\ndebloques', GREEN_MID),
    ('60%', 'Femmes\nbeneficiaires', BLUE),
    ('-40%', 'Reduction taux\nde defaut', GOLD),
]
for i, (val, label, color) in enumerate(impacts):
    add_kpi_card(slide, Inches(0.5 + i * 3.15), Inches(1.8), val, label, color)

# ODD alignment
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4), 'Alignement ODD', 14, True, DARK)
odds = [
    'ODD 1 — Acces au credit pour les plus vulnerables',
    'ODD 2 — Financement des intrants agricoles',
    'ODD 5 — Scoring inclusif (GIE feminin, tontines)',
    'ODD 8 — Formalisation de l\'activite economique',
    'ODD 10 — Inclusion financiere des ruraux exclus',
]
add_bullet_list(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(3), odds, 12, DARK)

# Right side — Why it works
add_text_box(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(0.4), 'Pourquoi ca fonctionne', 14, True, DARK)
reasons = [
    'Agents existants (ANCAR, cooperatives)',
    'Donnees deja la (tontines, Mobile Money, ventes)',
    'USSD pour les 70% sans smartphone',
    'Consentement en langue locale + temoin',
    'Scoring inclusif : pas besoin de titre foncier',
    'Boucle post-credit : le score s\'ameliore',
]
add_bullet_list(slide, Inches(7), Inches(4.2), Inches(5.5), Inches(3), reasons, 12, DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 : SECURITE & ANTI-FRAUDE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'CREDIBILITE & ANTI-FRAUDE', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'Les 5 risques critiques securises dans le prototype', 20, True, DARK)

risks = [
    ('Score trop binaire', 'Preuves alternatives : droit coutumier,\nGIE, tontine, validation communautaire', GREEN_DARK),
    ('Consentement insuffisant', 'Langue locale, lecture orale par agent,\npresence temoin, trace audit', GREEN_MID),
    ('Fraude agent', 'Score fiabilite agent, croisement donnees,\ntaux coherence, historique', BLUE),
    ('Credit non adapte', 'Recommandation calendrier saisonnier\naligne sur cycles recolte/vente', GOLD),
    ('Pas de suivi post-credit', 'Boucle remboursement/retard/defaut\nqui enrichit l\'historique', RGBColor(0x74, 0x52, 0x6F)),
]

for i, (risk, solution, color) in enumerate(risks):
    top = Inches(1.7 + i * 1.1)
    # Risk label
    shape = add_shape_card(slide, Inches(0.8), top, Inches(3.5), Inches(0.9), color, risk, 11, WHITE)
    # Arrow
    add_text_box(slide, Inches(4.5), top + Inches(0.2), Inches(0.5), Inches(0.5), '>', 18, True, GREEN_MID, PP_ALIGN.CENTER)
    # Solution
    add_shape_card(slide, Inches(5.2), top, Inches(7.3), Inches(0.9), GREEN_PALE, solution, 10, DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 : ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'FEUILLE DE ROUTE', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'Du hackathon a l\'echelle UEMOA en 36 mois', 20, True, DARK)

# Timeline bar
bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.3))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN_PALE
bar.line.fill.background()

phases = [
    ('Hackathon\nMai 2026', 'MVP valide\nScoring 13 criteres\nCollecte agent\nBoucle post-credit', Inches(0.8), GREEN_DARK),
    ('Pilote terrain\n6 mois', '100 agriculteurs\n5 agents formes\n2 IMF partenaires\n1er credit debloque', Inches(3.7), GREEN_MID),
    ('Croissance\n6-18 mois', '3 pays UEMOA\n25 agents\n5 000 agriculteurs\nAPI Mobile Money', Inches(6.6), BLUE),
    ('Echelle\n18-36 mois', '8 pays UEMOA\n80 agents\n25 000 agriculteurs\n30+ IMF connectees', Inches(9.5), GOLD),
]

for i, (title, desc, left, color) in enumerate(phases):
    # Dot on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.9), Inches(3.35), Inches(0.5), Inches(0.5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    tf = dot.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(i+1)
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title above
    add_text_box(slide, left, Inches(1.8), Inches(2.6), Inches(1.2), title, 12, True, color, PP_ALIGN.CENTER)

    # Description below
    add_text_box(slide, left, Inches(4.1), Inches(2.6), Inches(2.5), desc, 10, False, DARK, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 : EQUIPE & AVANTAGE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), 'NOTRE AVANTAGE', 12, True, GREEN_MID)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6), 'Pourquoi FresCoop va reussir la ou les autres echouent', 20, True, DARK)

advantages = [
    ('Realisme', 'Pas de techno-solutionnisme.\nNous partons des acteurs existants\n(agents, cooperatives, tontines).'),
    ('Focus', 'Nous ne faisons pas tout.\nNous faisons UNE chose bien :\nrendre l\'agriculteur financable.'),
    ('Terrain', 'Consentement en langue locale,\ntemoin, foncier coutumier.\nNous comprenons la realite.'),
    ('Modele viable', 'Les IMF paient parce qu\'on leur\nfait gagner du temps et\nreduire le risque credit.'),
    ('Scalable', 'Agent + USSD + API.\nFonctionne sans smartphone,\nmeme en zone 2G.'),
]

for i, (title, desc) in enumerate(advantages):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.1)
    top = Inches(1.7 + row * 2.8)

    card = add_shape_card(slide, left, top, Inches(3.8), Inches(2.4), GREEN_PALE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.4), Inches(0.4), title, 14, True, GREEN_DARK, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.4), Inches(1.5), desc, 11, False, DARK, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 : CLOSING / ASK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, GREEN_DARK, GREEN_MID)

add_text_box(slide, Inches(2), Inches(1.5), Inches(9), Inches(1), 'FresCoop', 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.6), 'De la recolte au credit — la chaine de valeur agricole digitalisee', 18, False, GREEN_PALE, PP_ALIGN.CENTER)

# Ask boxes
add_shape_card(slide, Inches(2), Inches(3.6), Inches(4.2), Inches(1.8), RGBColor(0x13, 0x6B, 0x5A))
add_text_box(slide, Inches(2.2), Inches(3.8), Inches(3.8), Inches(0.5), '15M FCFA', 28, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.2), Inches(4.4), Inches(3.8), Inches(0.8), 'Financement pre-seed\npour valider le pilote terrain', 11, False, GREEN_PALE, PP_ALIGN.CENTER)

add_shape_card(slide, Inches(7), Inches(3.6), Inches(4.2), Inches(1.8), RGBColor(0x13, 0x6B, 0x5A))
add_text_box(slide, Inches(7.2), Inches(3.8), Inches(3.8), Inches(0.5), '6 mois', 28, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(3.8), Inches(0.8), 'Pour prouver le modele :\n1er credit debloque via FresCoop', 11, False, GREEN_PALE, PP_ALIGN.CENTER)

# Final quote
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1), '"Notre ambition est claire : ne pas remplacer la banque,\nmais lui donner les preuves necessaires pour financer ceux qui nourrissent le pays."', 13, False, RGBColor(0xA7, 0xF3, 0xD0), PP_ALIGN.CENTER)

# Contact
add_text_box(slide, Inches(3), Inches(6.8), Inches(7), Inches(0.5), 'Equipe FresCoop  |  Hackathon GIM-UEMOA 2026  |  Friction 4', 10, False, RGBColor(0x6E, 0xE7, 0xB7), PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'FRESCOOP_PITCH_DECK.pptx')
prs.save(output_path)
print(f'Presentation sauvegardee : {output_path}')
print(f'Nombre de slides : {len(prs.slides)}')
