from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN_DARK = RGBColor(0x0A, 0x3D, 0x2A)
GREEN_700 = RGBColor(0x1F, 0x83, 0x5D)
GREEN_500 = RGBColor(0x10, 0xB9, 0x81)
GREEN_100 = RGBColor(0xD1, 0xFA, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_600 = RGBColor(0xDC, 0x26, 0x26)
AMBER_600 = RGBColor(0xD9, 0x77, 0x06)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
GRAY_600 = RGBColor(0x4B, 0x55, 0x63)
BLUE_600 = RGBColor(0x25, 0x63, 0xEB)
PURPLE_600 = RGBColor(0x7C, 0x3A, 0xED)
GOLD_600 = RGBColor(0xCA, 0x8A, 0x04)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_slide_gradient_bg(slide, color1, color2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_multi_text(slide, left, top, width, height, paragraphs_data):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pdata.get('text', '')
        p.font.size = Pt(pdata.get('size', 14))
        p.font.bold = pdata.get('bold', False)
        p.font.color.rgb = pdata.get('color', WHITE)
        p.font.name = pdata.get('font', 'Segoe UI')
        p.alignment = pdata.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(pdata.get('space_after', 6))
    return txBox

# ============================================================================
# SLIDE 1: COVER
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, GREEN_DARK)

# Accent bar top
bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN_500
bar.line.fill.background()

# Hackathon badge
add_text_box(slide1, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             'HACKATHON FILIÈRES AGRICOLES GIM-UEMOA 2026 — POINT 4 : ACCÈS AU FINANCEMENT',
             font_size=11, bold=True, color=GREEN_500)

# Main title
add_text_box(slide1, Inches(0.8), Inches(2.0), Inches(8), Inches(1.2),
             'De l\'invisible au finançable.',
             font_size=44, bold=True, color=WHITE, font_name='Segoe UI')

add_text_box(slide1, Inches(0.8), Inches(3.1), Inches(6), Inches(0.9),
             'En 90 jours.',
             font_size=44, bold=True, color=GREEN_500, font_name='Segoe UI')

# Subtitle
add_text_box(slide1, Inches(0.8), Inches(4.2), Inches(7.5), Inches(1.2),
             'FresCoop transforme chaque vente, chaque livraison, chaque paiement d\'un agriculteur en preuve bancaire vérifiable — construisant automatiquement un score de crédit exploitable par toute institution financière de la zone UEMOA.',
             font_size=16, color=RGBColor(0xD1, 0xD5, 0xDB))

# Stats row
stats = [('80%', 'des agriculteurs\nexclus du crédit'), ('60M', "d'agriculteurs\ndans l'UEMOA"), ('3%', "du crédit va\nà l'agriculture")]
for i, (val, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.2)
    rect = add_rounded_rect(slide1, x, Inches(5.8), Inches(2.8), Inches(1.3), RGBColor(0x14, 0x4D, 0x36))
    add_text_box(slide1, x + Inches(0.3), Inches(5.9), Inches(2.4), Inches(0.6), val, font_size=28, bold=True, color=GREEN_500)
    add_text_box(slide1, x + Inches(0.3), Inches(6.4), Inches(2.4), Inches(0.7), label, font_size=11, color=RGBColor(0xA1, 0xA1, 0xAA))

# FresCoop logo
logo_shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.5), Inches(0.7), Inches(0.7))
logo_shape.fill.solid()
logo_shape.fill.fore_color.rgb = GREEN_500
logo_shape.line.fill.background()
add_text_box(slide1, Inches(11.55), Inches(0.55), Inches(0.6), Inches(0.6), 'F', font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide1, Inches(12.2), Inches(0.55), Inches(1), Inches(0.5), 'FresCoop', font_size=14, bold=True, color=WHITE)

# ============================================================================
# SLIDE 2: PROBLÈME
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, WHITE)

# Section header
add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
             'PROBLÈME', font_size=11, bold=True, color=GREEN_700)

add_text_box(slide2, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             'Les agriculteurs sont invisibles pour le système financier.',
             font_size=32, bold=True, color=SLATE_800, font_name='Segoe UI')

add_text_box(slide2, Inches(0.8), Inches(1.7), Inches(9), Inches(0.6),
             'Pas parce qu\'ils ne sont pas fiables — parce qu\'ils n\'ont aucune preuve exploitable de leur activité économique.',
             font_size=14, color=SLATE_500)

# Three cards
cards = [
    ('80%', 'Sans accès au crédit', 'Les institutions financières ne financent pas ceux qu\'elles ne peuvent pas évaluer.\nPas de relevé bancaire = pas de crédit.', RED_600),
    ('3%', 'Du crédit vers l\'agriculture', 'Alors que l\'agriculture représente 35% du PIB de la zone UEMOA.\nUn paradoxe inacceptable.', AMBER_600),
    ('0', 'Historique financier exploitable', 'Chaque jour, des millions de transactions en espèces ne laissent aucune trace pour les banques.', SLATE_800),
]

for i, (stat, title, desc, accent) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.8)

    # Card background
    card = add_rounded_rect(slide2, x, y, Inches(3.7), Inches(4.2), RGBColor(0xF9, 0xFA, 0xFB))

    # Top accent bar
    accent_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.7), Inches(0.06))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    # Stat
    add_text_box(slide2, x + Inches(0.4), y + Inches(0.5), Inches(3), Inches(0.9), stat, font_size=48, bold=True, color=accent)

    # Title
    add_text_box(slide2, x + Inches(0.4), y + Inches(1.5), Inches(3), Inches(0.7), title, font_size=13, bold=True, color=SLATE_800)

    # Description
    add_text_box(slide2, x + Inches(0.4), y + Inches(2.2), Inches(3), Inches(1.8), desc, font_size=12, color=GRAY_600)

# ============================================================================
# SLIDE 3: SOLUTION
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, WHITE)

add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
             'SOLUTION', font_size=11, bold=True, color=GREEN_700)

add_text_box(slide3, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             'Chaque transaction devient une preuve. Chaque preuve construit un score.',
             font_size=30, bold=True, color=SLATE_800)

add_text_box(slide3, Inches(0.8), Inches(1.7), Inches(9), Inches(0.6),
             'FresCoop capture automatiquement l\'activité économique réelle et la transforme en dossier bancaire portable — sans changer les habitudes de l\'agriculteur.',
             font_size=14, color=SLATE_500)

# Solution steps
solutions = [
    ('1', 'Captation automatique', 'Chaque vente, livraison, paiement génère des données de scoring sans effort.'),
    ('2', 'Scoring intelligent', 'Score de 0 à 100 calculé en temps réel. Pas de déclarations, des preuves.'),
    ('3', 'Dossier bancaire portable', 'PDF vérifiable par QR code, présentable à toute banque ou SFD.'),
    ('4', 'Accès au crédit', 'En 3 mois d\'activité, passez d\'invisible à finançable. Demande en 1 clic.'),
    ('5', 'Vérification terrain', 'Des agents valident l\'identité et accompagnent la progression.'),
]

for i, (num, title, desc) in enumerate(solutions):
    x = Inches(0.8 + i * 2.45)
    y = Inches(2.8)

    # Card
    card = add_rounded_rect(slide3, x, y, Inches(2.25), Inches(4.2), RGBColor(0xF0, 0xFD, 0xF4))

    # Number circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.4), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_700
    circle.line.fill.background()
    add_text_box(slide3, x + Inches(0.35), y + Inches(0.42), Inches(0.4), Inches(0.45), num, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide3, x + Inches(0.3), y + Inches(1.2), Inches(1.8), Inches(0.8), title, font_size=13, bold=True, color=GREEN_DARK)

    # Desc
    add_text_box(slide3, x + Inches(0.3), y + Inches(2.1), Inches(1.8), Inches(1.8), desc, font_size=11, color=GRAY_600)

# ============================================================================
# SLIDE 4: MODÈLE ÉCONOMIQUE
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, WHITE)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
             'MODELE ECONOMIQUE', font_size=11, bold=True, color=GREEN_700)

add_text_box(slide4, Inches(0.8), Inches(0.8), Inches(10), Inches(0.7),
             'Gratuit pour l\'agriculteur. Rentable pour tous.',
             font_size=30, bold=True, color=SLATE_800)

# Flow diagram: Qui paie ?
flow_bg = add_rounded_rect(slide4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.5), RGBColor(0xF0, 0xFD, 0xF4))
add_text_box(slide4, Inches(1.2), Inches(1.7), Inches(2), Inches(0.3),
             'QUI PAIE ?', font_size=9, bold=True, color=GREEN_700)

flow_actors = [
    ('Agriculteur', 'GRATUIT'),
    ('Acheteur', 'Paie le produit'),
    ('SFD / Banque', 'Paie l\'acces au score'),
]
for i, (actor, role) in enumerate(flow_actors):
    x = Inches(1.5 + i * 4.0)
    actor_bg = add_rounded_rect(slide4, x, Inches(2.05), Inches(2.8), Inches(0.85), WHITE)
    add_text_box(slide4, x + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.4), actor, font_size=13, bold=True, color=SLATE_800)
    add_text_box(slide4, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.3), role, font_size=10, bold=True, color=GREEN_700)
    if i < 2:
        add_text_box(slide4, Inches(4.5 + i * 4.0), Inches(2.25), Inches(0.5), Inches(0.4), '→', font_size=20, color=GREEN_700, alignment=PP_ALIGN.CENTER)

# Revenue cards (compact)
revenues = [
    ('Commission marketplace', '60% du CA', '2-5% par transaction', 'L\'agriculteur vend plus cher, la commission est invisible dans le gain.'),
    ('Scoring-as-a-Service', '30% du CA', '150K-500K FCFA/mois par SFD', 'Les SFD economisent 80% vs l\'enquete terrain traditionnelle.'),
    ('Services a valeur ajoutee', '10% du CA', 'Commission 5-15%', 'Assurance recolte, intrants a credit, logistique — marches inaccessibles avant.'),
]

for i, (title, pct, pricing, desc) in enumerate(revenues):
    x = Inches(0.8 + i * 4.0)
    y = Inches(3.4)

    card = add_rounded_rect(slide4, x, y, Inches(3.7), Inches(2.6), RGBColor(0xF9, 0xFA, 0xFB))
    bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.7), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_700
    bar.line.fill.background()

    # % badge
    add_text_box(slide4, x + Inches(0.3), y + Inches(0.25), Inches(1.2), Inches(0.35), pct, font_size=10, bold=True, color=GREEN_700)

    # Title
    add_text_box(slide4, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Inches(0.5), title, font_size=14, bold=True, color=SLATE_800)

    # Pricing
    add_text_box(slide4, x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(0.35), pricing, font_size=10, bold=True, color=GREEN_700)

    # Description
    add_text_box(slide4, x + Inches(0.3), y + Inches(1.5), Inches(3.2), Inches(1.0), desc, font_size=10, color=GRAY_600)

# Viability metrics at bottom
viability = [
    ('Point mort', '500 agri + 3 SFD', '6 mois'),
    ('Cout marginal', '~0 FCFA/agri', '100% cloud'),
    ('TAM', '60M x 8 pays', 'Zone UEMOA'),
    ('Projection An 1', '5 000 utilisateurs', 'CA: 45M FCFA'),
]
for i, (label, val, sub) in enumerate(viability):
    x = Inches(0.8 + i * 3.1)
    y = Inches(6.2)
    rect = add_rounded_rect(slide4, x, y, Inches(2.8), Inches(0.95), RGBColor(0xF0, 0xFD, 0xF4))
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.08), Inches(2.4), Inches(0.3), label, font_size=9, bold=True, color=SLATE_500)
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.35), Inches(2.4), Inches(0.35), val, font_size=13, bold=True, color=GREEN_700)
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.68), Inches(2.4), Inches(0.25), sub, font_size=9, color=SLATE_500)

# ============================================================================
# SLIDE 5: IMPACT
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, WHITE)

add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             'IMPACT ÉCONOMIQUE & SOCIAL', font_size=11, bold=True, color=GREEN_700)

add_text_box(slide5, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             'Des résultats concrets. Mesurables. Transformateurs.',
             font_size=32, bold=True, color=SLATE_800)

# Economic impact metrics
eco_metrics = [
    ('65%', 'deviennent bancables\nen 90 jours', GREEN_700),
    ('-80%', 'coût d\'évaluation\npour les SFD', BLUE_600),
    ('90j', 'du premier crédit\nobtenu via FresCoop', GOLD_600),
    ('×5', 'plus d\'agriculteurs\nfinancés par SFD', PURPLE_600),
]

for i, (val, label, color) in enumerate(eco_metrics):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.2)

    rect = add_rounded_rect(slide5, x, y, Inches(2.8), Inches(1.6), RGBColor(0xF9, 0xFA, 0xFB))
    add_text_box(slide5, x + Inches(0.3), y + Inches(0.2), Inches(2.4), Inches(0.7), val, font_size=36, bold=True, color=color)
    add_text_box(slide5, x + Inches(0.3), y + Inches(0.95), Inches(2.4), Inches(0.6), label, font_size=11, color=GRAY_600)

# Social impact
add_text_box(slide5, Inches(0.8), Inches(4.2), Inches(6), Inches(0.4),
             'IMPACT SUR L\'ACCÈS AU FINANCEMENT', font_size=11, bold=True, color=GREEN_700)

social_items = [
    ('Inclusion des femmes', '40% des agricultrices exclues accèdent enfin au crédit grâce à un scoring objectif sans biais de genre.'),
    ('Dossier bancaire portable', 'L\'agriculteur présente son score vérifié à n\'importe quelle banque ou SFD de la zone UEMOA.'),
    ('Crédit sans garantie foncière', 'Le scoring remplace l\'exigence de titre foncier — la preuve d\'activité suffit pour le financement.'),
    ('Réduction du risque SFD', 'Les SFD financent en connaissance de cause : données vérifiées, historique complet, défaut prévisible.'),
]

for i, (title, desc) in enumerate(social_items):
    x = Inches(0.8 + i * 3.1)
    y = Inches(4.7)

    rect = add_rounded_rect(slide5, x, y, Inches(2.8), Inches(2.4), RGBColor(0xF0, 0xFD, 0xF4))
    add_text_box(slide5, x + Inches(0.3), y + Inches(0.3), Inches(2.4), Inches(0.5), title, font_size=13, bold=True, color=GREEN_DARK)
    add_text_box(slide5, x + Inches(0.3), y + Inches(0.9), Inches(2.3), Inches(1.4), desc, font_size=11, color=GRAY_600)

# ============================================================================
# SLIDE 6: CLOSING / CTA
# ============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, GREEN_DARK)

# Accent bar
bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN_500
bar.line.fill.background()

add_text_box(slide6, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
             'CE N\'EST PAS UN PROTOTYPE. C\'EST UN PRODUIT FONCTIONNEL.',
             font_size=12, bold=True, color=GREEN_500)

add_text_box(slide6, Inches(0.8), Inches(2.3), Inches(10), Inches(1.2),
             'FresCoop rend chaque agriculteur\nvisible, vérifiable et finançable.',
             font_size=40, bold=True, color=WHITE)

add_text_box(slide6, Inches(0.8), Inches(3.8), Inches(8), Inches(0.8),
             'Grâce à la preuve économique portable — un dossier bancaire construit\nautomatiquement à partir de l\'activité réelle sur la plateforme.',
             font_size=16, color=RGBColor(0xD1, 0xD5, 0xDB))

# Demo URL
demo_rect = add_rounded_rect(slide6, Inches(0.8), Inches(5.2), Inches(5), Inches(0.7), GREEN_700)
add_text_box(slide6, Inches(1.0), Inches(5.25), Inches(4.6), Inches(0.6),
             '▶  Démo live : frescoopuemoa.up.railway.app',
             font_size=14, bold=True, color=WHITE)

# Contact
add_text_box(slide6, Inches(0.8), Inches(6.3), Inches(5), Inches(0.5),
             'Contact : contact@frescoop.sn',
             font_size=13, color=RGBColor(0x9C, 0xA3, 0xAF))

add_text_box(slide6, Inches(0.8), Inches(6.8), Inches(8), Inches(0.4),
             'Hackathon Filières Agricoles GIM-UEMOA 2026 — Point 4 : Accès au financement agricole',
             font_size=11, color=RGBColor(0x6B, 0x72, 0x80))

# Logo bottom right
logo2 = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(6.2), Inches(0.8), Inches(0.8))
logo2.fill.solid()
logo2.fill.fore_color.rgb = GREEN_500
logo2.line.fill.background()
add_text_box(slide6, Inches(11.56), Inches(6.28), Inches(0.68), Inches(0.68), 'F', font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'FRESCOOP_PITCH_v2.pptx')
prs.save(output_path)
print(f'Pitch deck saved: {output_path}')
